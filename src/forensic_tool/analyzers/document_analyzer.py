"""
Analisador especializado para documentos
"""

from pathlib import Path
from typing import Dict, Any, Optional, List
import logging
from datetime import datetime
from .base import MultiFormatAnalyzer

logger = logging.getLogger(__name__)

# Importações condicionais
try:
    import PyPDF2
    from PyPDF2 import PdfReader
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    import docx
    from docx.opc.coreprops import CoreProperties
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False

try:
    import pptx
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class DocumentAnalyzer(MultiFormatAnalyzer):
    """Analisador para documentos de escritório"""
    
    def __init__(self):
        super().__init__("DocumentAnalyzer")
        
        # Registrar handlers para diferentes formatos
        if PDF_AVAILABLE:
            self.add_format_handler({'.pdf'}, self._analyze_pdf)
        
        if DOCX_AVAILABLE:
            self.add_format_handler({'.docx', '.doc'}, self._analyze_word)
        
        if EXCEL_AVAILABLE:
            self.add_format_handler({'.xlsx', '.xls'}, self._analyze_excel)
        
        if PPTX_AVAILABLE:
            self.add_format_handler({'.pptx', '.ppt'}, self._analyze_powerpoint)
        
        # Texto simples sempre disponível
        self.add_format_handler({'.txt', '.rtf', '.csv'}, self._analyze_text)
        
        # Log de disponibilidade
        if not PDF_AVAILABLE:
            logger.warning("PyPDF2 não disponível - análise de PDF desabilitada")
        if not DOCX_AVAILABLE:
            logger.warning("python-docx não disponível - análise de Word desabilitada")
        if not EXCEL_AVAILABLE:
            logger.warning("openpyxl não disponível - análise de Excel desabilitada")
        if not PPTX_AVAILABLE:
            logger.warning("python-pptx não disponível - análise de PowerPoint desabilitada")
    
    def _analyze_pdf(self, file_path: Path) -> Dict[str, Any]:
        """Análise específica para arquivos PDF"""
        metadata = {
            'document_type': 'PDF',
            'pages': 0,
            'encrypted': False,
            'metadata': {},
            'text_content': {},
            'security_info': {},
            'structure_info': {}
        }
        
        try:
            with open(file_path, 'rb') as file:
                reader = PdfReader(file)
                
                # Informações básicas
                metadata['pages'] = len(reader.pages)
                metadata['encrypted'] = reader.is_encrypted
                
                # Metadados do documento
                if reader.metadata:
                    doc_metadata = {}
                    for key, value in reader.metadata.items():
                        clean_key = key.replace('/', '') if key.startswith('/') else key
                        doc_metadata[clean_key] = str(value) if value else None
                    metadata['metadata'] = doc_metadata
                
                # Análise de conteúdo (primeiras páginas)
                text_analysis = self._analyze_pdf_content(reader)
                metadata['text_content'] = text_analysis
                
                # Informações de segurança
                metadata['security_info'] = self._analyze_pdf_security(reader)
                
                # Estrutura do documento
                metadata['structure_info'] = self._analyze_pdf_structure(reader)
                
        except Exception as e:
            logger.error(f"Erro na análise PDF {file_path}: {e}")
            metadata['analysis_error'] = str(e)
        
        return metadata
    
    def _analyze_pdf_content(self, reader: PdfReader, max_pages: int = 5) -> Dict[str, Any]:
        """Análise do conteúdo textual do PDF"""
        content_info = {
            'total_characters': 0,
            'total_words': 0,
            'sample_text': '',
            'languages_detected': [],
            'has_text': False
        }
        
        try:
            all_text = ""
            pages_analyzed = min(len(reader.pages), max_pages)
            
            for i in range(pages_analyzed):
                try:
                    page = reader.pages[i]
                    page_text = page.extract_text()
                    if page_text:
                        all_text += page_text + "\n"
                except Exception as e:
                    logger.debug(f"Erro ao extrair texto da página {i}: {e}")
            
            if all_text.strip():
                content_info['has_text'] = True
                content_info['total_characters'] = len(all_text)
                content_info['total_words'] = len(all_text.split())
                content_info['sample_text'] = all_text[:500] + "..." if len(all_text) > 500 else all_text
                
                # Detecção básica de idioma (simplificada)
                content_info['languages_detected'] = self._detect_languages(all_text)
            
        except Exception as e:
            logger.debug(f"Erro na análise de conteúdo PDF: {e}")
            content_info['content_analysis_error'] = str(e)
        
        return content_info
    
    def _analyze_pdf_security(self, reader: PdfReader) -> Dict[str, Any]:
        """Análise de segurança do PDF"""
        security_info = {
            'encrypted': reader.is_encrypted,
            'permissions': {},
            'security_handler': 'unknown'
        }
        
        try:
            if hasattr(reader, 'decrypt_when_encrypted'):
                # Tentar obter informações de permissões
                if hasattr(reader, 'trailer') and reader.trailer:
                    encrypt_dict = reader.trailer.get('/Encrypt')
                    if encrypt_dict:
                        security_info['security_handler'] = str(encrypt_dict.get('/Filter', 'unknown'))
                        
                        # Permissões (se disponível)
                        permissions = encrypt_dict.get('/P')
                        if permissions:
                            security_info['permissions'] = self._decode_pdf_permissions(permissions)
            
        except Exception as e:
            logger.debug(f"Erro na análise de segurança PDF: {e}")
            security_info['security_analysis_error'] = str(e)
        
        return security_info
    
    def _analyze_pdf_structure(self, reader: PdfReader) -> Dict[str, Any]:
        """Análise da estrutura do PDF"""
        structure_info = {
            'has_bookmarks': False,
            'has_forms': False,
            'has_annotations': False,
            'pdf_version': 'unknown'
        }
        
        try:
            # Versão do PDF
            if hasattr(reader, 'pdf_header'):
                structure_info['pdf_version'] = reader.pdf_header
            
            # Bookmarks/Outlines
            if hasattr(reader, 'outline') and reader.outline:
                structure_info['has_bookmarks'] = True
                structure_info['bookmark_count'] = len(reader.outline)
            
            # Formulários e anotações (análise básica)
            for page in reader.pages[:5]:  # Verificar primeiras 5 páginas
                if hasattr(page, 'annotations') and page.annotations:
                    structure_info['has_annotations'] = True
                
                # Verificar se há campos de formulário
                if '/AcroForm' in str(page):
                    structure_info['has_forms'] = True
            
        except Exception as e:
            logger.debug(f"Erro na análise de estrutura PDF: {e}")
            structure_info['structure_analysis_error'] = str(e)
        
        return structure_info
    
    def _decode_pdf_permissions(self, permissions: int) -> Dict[str, bool]:
        """Decodifica permissões do PDF"""
        # Baseado na especificação PDF
        return {
            'print': bool(permissions & 4),
            'modify': bool(permissions & 8),
            'copy': bool(permissions & 16),
            'add_annotations': bool(permissions & 32),
            'fill_forms': bool(permissions & 256),
            'extract_for_accessibility': bool(permissions & 512),
            'assemble': bool(permissions & 1024),
            'print_high_quality': bool(permissions & 2048)
        }
    
    def _analyze_word(self, file_path: Path) -> Dict[str, Any]:
        """Análise específica para documentos Word"""
        metadata = {
            'document_type': 'Word Document',
            'properties': {},
            'content_stats': {},
            'structure_info': {},
            'revision_info': {}
        }
        
        try:
            doc = docx.Document(file_path)
            
            # Propriedades do documento
            if doc.core_properties:
                props = doc.core_properties
                metadata['properties'] = {
                    'title': props.title,
                    'author': props.author,
                    'subject': props.subject,
                    'keywords': props.keywords,
                    'comments': props.comments,
                    'created': props.created.isoformat() if props.created else None,
                    'modified': props.modified.isoformat() if props.modified else None,
                    'last_modified_by': props.last_modified_by,
                    'revision': props.revision,
                    'category': props.category,
                    'language': props.language
                }
            
            # Estatísticas de conteúdo
            metadata['content_stats'] = self._analyze_word_content(doc)
            
            # Estrutura do documento
            metadata['structure_info'] = self._analyze_word_structure(doc)
            
        except Exception as e:
            logger.error(f"Erro na análise Word {file_path}: {e}")
            metadata['analysis_error'] = str(e)
        
        return metadata
    
    def _analyze_word_content(self, doc) -> Dict[str, Any]:
        """Análise do conteúdo do documento Word"""
        content_stats = {
            'paragraphs': len(doc.paragraphs),
            'total_characters': 0,
            'total_words': 0,
            'has_tables': False,
            'has_images': False,
            'sample_text': ''
        }
        
        try:
            # Contar texto
            all_text = ""
            for paragraph in doc.paragraphs:
                if paragraph.text:
                    all_text += paragraph.text + "\n"
            
            content_stats['total_characters'] = len(all_text)
            content_stats['total_words'] = len(all_text.split())
            content_stats['sample_text'] = all_text[:500] + "..." if len(all_text) > 500 else all_text
            
            # Verificar tabelas
            if doc.tables:
                content_stats['has_tables'] = True
                content_stats['table_count'] = len(doc.tables)
            
            # Verificar imagens (através de relacionamentos)
            if hasattr(doc, 'part') and hasattr(doc.part, 'related_parts'):
                for rel in doc.part.related_parts.values():
                    if 'image' in str(rel.content_type).lower():
                        content_stats['has_images'] = True
                        break
            
        except Exception as e:
            logger.debug(f"Erro na análise de conteúdo Word: {e}")
            content_stats['content_analysis_error'] = str(e)
        
        return content_stats
    
    def _analyze_word_structure(self, doc) -> Dict[str, Any]:
        """Análise da estrutura do documento Word"""
        structure_info = {
            'sections': 0,
            'headers_footers': False,
            'footnotes': False,
            'styles_used': []
        }
        
        try:
            # Seções
            structure_info['sections'] = len(doc.sections)
            
            # Headers e footers
            for section in doc.sections:
                if section.header.paragraphs or section.footer.paragraphs:
                    structure_info['headers_footers'] = True
                    break
            
            # Estilos utilizados (simplificado)
            styles_used = set()
            for paragraph in doc.paragraphs[:50]:  # Primeiros 50 parágrafos
                if paragraph.style and paragraph.style.name:
                    styles_used.add(paragraph.style.name)
            
            structure_info['styles_used'] = list(styles_used)[:10]  # Top 10
            
        except Exception as e:
            logger.debug(f"Erro na análise de estrutura Word: {e}")
            structure_info['structure_analysis_error'] = str(e)
        
        return structure_info
    
    def _analyze_excel(self, file_path: Path) -> Dict[str, Any]:
        """Análise específica para planilhas Excel"""
        metadata = {
            'document_type': 'Excel Spreadsheet',
            'workbook_info': {},
            'worksheets': [],
            'data_analysis': {},
            'formulas_info': {}
        }
        
        try:
            workbook = openpyxl.load_workbook(file_path, read_only=True, data_only=False)
            
            # Informações do workbook
            metadata['workbook_info'] = {
                'worksheet_count': len(workbook.worksheets),
                'worksheet_names': workbook.sheetnames,
                'active_sheet': workbook.active.title if workbook.active else None
            }
            
            # Análise das planilhas
            for sheet in workbook.worksheets:
                sheet_info = self._analyze_excel_worksheet(sheet)
                metadata['worksheets'].append(sheet_info)
            
            # Análise geral de dados
            metadata['data_analysis'] = self._analyze_excel_data(workbook)
            
        except Exception as e:
            logger.error(f"Erro na análise Excel {file_path}: {e}")
            metadata['analysis_error'] = str(e)
        
        return metadata
    
    def _analyze_excel_worksheet(self, sheet) -> Dict[str, Any]:
        """Análise de uma planilha específica"""
        sheet_info = {
            'name': sheet.title,
            'dimensions': f"{sheet.max_row}x{sheet.max_column}",
            'max_row': sheet.max_row,
            'max_column': sheet.max_column,
            'has_data': False,
            'cell_count': 0,
            'formula_count': 0
        }
        
        try:
            # Contar células com dados e fórmulas
            data_cells = 0
            formula_cells = 0
            
            # Limitar análise para performance
            max_rows = min(sheet.max_row, 1000)
            max_cols = min(sheet.max_column, 100)
            
            for row in sheet.iter_rows(max_row=max_rows, max_col=max_cols):
                for cell in row:
                    if cell.value is not None:
                        data_cells += 1
                        if isinstance(cell.value, str) and cell.value.startswith('='):
                            formula_cells += 1
            
            sheet_info['has_data'] = data_cells > 0
            sheet_info['cell_count'] = data_cells
            sheet_info['formula_count'] = formula_cells
            
        except Exception as e:
            logger.debug(f"Erro na análise da planilha {sheet.title}: {e}")
            sheet_info['analysis_error'] = str(e)
        
        return sheet_info
    
    def _analyze_excel_data(self, workbook) -> Dict[str, Any]:
        """Análise geral dos dados do Excel"""
        data_analysis = {
            'total_cells_with_data': 0,
            'total_formulas': 0,
            'data_types_found': set(),
            'has_charts': False,
            'has_pivot_tables': False
        }
        
        try:
            for sheet in workbook.worksheets:
                # Análise limitada para performance
                max_rows = min(sheet.max_row, 100)
                max_cols = min(sheet.max_column, 50)
                
                for row in sheet.iter_rows(max_row=max_rows, max_col=max_cols):
                    for cell in row:
                        if cell.value is not None:
                            data_analysis['total_cells_with_data'] += 1
                            
                            # Tipo de dados
                            data_type = type(cell.value).__name__
                            data_analysis['data_types_found'].add(data_type)
                            
                            # Fórmulas
                            if isinstance(cell.value, str) and cell.value.startswith('='):
                                data_analysis['total_formulas'] += 1
            
            # Converter set para list para serialização
            data_analysis['data_types_found'] = list(data_analysis['data_types_found'])
            
        except Exception as e:
            logger.debug(f"Erro na análise de dados Excel: {e}")
            data_analysis['data_analysis_error'] = str(e)
        
        return data_analysis
    
    def _analyze_powerpoint(self, file_path: Path) -> Dict[str, Any]:
        """Análise específica para apresentações PowerPoint"""
        metadata = {
            'document_type': 'PowerPoint Presentation',
            'presentation_info': {},
            'slides': [],
            'content_analysis': {}
        }
        
        try:
            presentation = pptx.Presentation(file_path)
            
            # Informações da apresentação
            metadata['presentation_info'] = {
                'slide_count': len(presentation.slides),
                'slide_layouts_count': len(presentation.slide_layouts),
                'slide_masters_count': len(presentation.slide_masters)
            }
            
            # Análise dos slides
            for i, slide in enumerate(presentation.slides):
                slide_info = self._analyze_powerpoint_slide(slide, i)
                metadata['slides'].append(slide_info)
            
            # Análise geral do conteúdo
            metadata['content_analysis'] = self._analyze_powerpoint_content(presentation)
            
        except Exception as e:
            logger.error(f"Erro na análise PowerPoint {file_path}: {e}")
            metadata['analysis_error'] = str(e)
        
        return metadata
    
    def _analyze_powerpoint_slide(self, slide, slide_number: int) -> Dict[str, Any]:
        """Análise de um slide específico"""
        slide_info = {
            'slide_number': slide_number + 1,
            'layout_name': slide.slide_layout.name if slide.slide_layout else 'Unknown',
            'shape_count': len(slide.shapes),
            'has_text': False,
            'has_images': False,
            'has_tables': False,
            'text_content': ''
        }
        
        try:
            all_text = ""
            
            for shape in slide.shapes:
                # Texto
                if hasattr(shape, 'text') and shape.text:
                    all_text += shape.text + "\n"
                    slide_info['has_text'] = True
                
                # Imagens
                if hasattr(shape, 'image'):
                    slide_info['has_images'] = True
                
                # Tabelas
                if hasattr(shape, 'table'):
                    slide_info['has_tables'] = True
            
            slide_info['text_content'] = all_text[:300] + "..." if len(all_text) > 300 else all_text
            
        except Exception as e:
            logger.debug(f"Erro na análise do slide {slide_number}: {e}")
            slide_info['analysis_error'] = str(e)
        
        return slide_info
    
    def _analyze_powerpoint_content(self, presentation) -> Dict[str, Any]:
        """Análise geral do conteúdo da apresentação"""
        content_analysis = {
            'total_text_characters': 0,
            'total_shapes': 0,
            'slides_with_images': 0,
            'slides_with_tables': 0,
            'common_layouts': {}
        }
        
        try:
            layout_count = {}
            
            for slide in presentation.slides:
                content_analysis['total_shapes'] += len(slide.shapes)
                
                # Layout
                layout_name = slide.slide_layout.name if slide.slide_layout else 'Unknown'
                layout_count[layout_name] = layout_count.get(layout_name, 0) + 1
                
                # Conteúdo
                has_images = False
                has_tables = False
                
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text:
                        content_analysis['total_text_characters'] += len(shape.text)
                    
                    if hasattr(shape, 'image'):
                        has_images = True
                    
                    if hasattr(shape, 'table'):
                        has_tables = True
                
                if has_images:
                    content_analysis['slides_with_images'] += 1
                if has_tables:
                    content_analysis['slides_with_tables'] += 1
            
            content_analysis['common_layouts'] = layout_count
            
        except Exception as e:
            logger.debug(f"Erro na análise de conteúdo PowerPoint: {e}")
            content_analysis['content_analysis_error'] = str(e)
        
        return content_analysis
    
    def _analyze_text(self, file_path: Path) -> Dict[str, Any]:
        """Análise para arquivos de texto simples"""
        metadata = {
            'document_type': 'Text Document',
            'encoding': 'unknown',
            'content_stats': {},
            'text_analysis': {}
        }
        
        try:
            # Detectar encoding
            encodings_to_try = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
            content = ""
            used_encoding = 'unknown'
            
            for encoding in encodings_to_try:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        content = f.read()
                    used_encoding = encoding
                    break
                except UnicodeDecodeError:
                    continue
            
            metadata['encoding'] = used_encoding
            
            if content:
                # Estatísticas básicas
                lines = content.splitlines()
                words = content.split()
                
                metadata['content_stats'] = {
                    'total_characters': len(content),
                    'total_lines': len(lines),
                    'total_words': len(words),
                    'average_line_length': sum(len(line) for line in lines) / len(lines) if lines else 0,
                    'empty_lines': sum(1 for line in lines if not line.strip()),
                    'max_line_length': max(len(line) for line in lines) if lines else 0
                }
                
                # Análise de texto
                metadata['text_analysis'] = {
                    'sample_text': content[:500] + "..." if len(content) > 500 else content,
                    'languages_detected': self._detect_languages(content),
                    'has_urls': 'http' in content.lower() or 'www.' in content.lower(),
                    'has_emails': '@' in content and '.' in content
                }
            
        except Exception as e:
            logger.error(f"Erro na análise de texto {file_path}: {e}")
            metadata['analysis_error'] = str(e)
        
        return metadata
    
    def _detect_languages(self, text: str) -> List[str]:
        """Detecção básica de idiomas (simplificada)"""
        # Implementação muito básica - pode ser melhorada com bibliotecas especializadas
        languages = []
        
        try:
            text_lower = text.lower()
            
            # Palavras comuns em diferentes idiomas
            language_indicators = {
                'portuguese': ['que', 'para', 'com', 'uma', 'por', 'não', 'são', 'dos', 'mais'],
                'english': ['the', 'and', 'for', 'are', 'but', 'not', 'you', 'all', 'can'],
                'spanish': ['que', 'para', 'con', 'una', 'por', 'son', 'los', 'más', 'como'],
                'french': ['que', 'pour', 'avec', 'une', 'par', 'sont', 'les', 'plus', 'comme']
            }
            
            for language, indicators in language_indicators.items():
                score = sum(1 for word in indicators if word in text_lower)
                if score >= 3:  # Threshold arbitrário
                    languages.append(language)
            
            if not languages:
                languages.append('unknown')
                
        except Exception as e:
            logger.debug(f"Erro na detecção de idiomas: {e}")
            languages = ['unknown']
        
        return languages
